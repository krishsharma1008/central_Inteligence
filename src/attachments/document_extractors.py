"""
Document extractors for various file formats.
Supports PDF, DOCX, XLSX, PPTX, and text files.
"""
import logging
from typing import Dict, Any, Optional
from io import BytesIO

logger = logging.getLogger('outlook-email.extractors')


class PDFExtractor:
    """Extract text from PDF files using PyMuPDF (text-based PDFs only)."""

    def extract(self, pdf_bytes: bytes) -> Dict[str, Any]:
        """
        Extract text from PDF.

        Args:
            pdf_bytes: PDF file content as bytes

        Returns:
            Dictionary with:
                - text: Extracted text content
                - page_count: Number of pages
                - metadata: PDF metadata
        """
        try:
            import fitz  # PyMuPDF

            # Open PDF from bytes
            pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")

            # Extract text from all pages
            text_parts = []
            for page_num in range(pdf_document.page_count):
                page = pdf_document[page_num]
                text = page.get_text()
                if text.strip():
                    text_parts.append(f"[Page {page_num + 1}]\n{text}")

            full_text = "\n\n".join(text_parts)

            # Get metadata
            metadata = {
                'page_count': pdf_document.page_count,
                'author': pdf_document.metadata.get('author', ''),
                'title': pdf_document.metadata.get('title', ''),
                'subject': pdf_document.metadata.get('subject', ''),
                'creator': pdf_document.metadata.get('creator', '')
            }

            pdf_document.close()

            logger.info(f"Extracted {len(full_text)} characters from {metadata['page_count']}-page PDF")

            return {
                'text': full_text,
                'page_count': metadata['page_count'],
                'metadata': metadata
            }

        except Exception as e:
            logger.error(f"Error extracting PDF: {str(e)}", exc_info=True)
            return {
                'text': '',
                'page_count': 0,
                'metadata': {},
                'error': str(e)
            }


class DOCXExtractor:
    """Extract text from Word documents using python-docx."""

    def extract(self, docx_bytes: bytes) -> Dict[str, Any]:
        """
        Extract text from DOCX file.

        Args:
            docx_bytes: DOCX file content as bytes

        Returns:
            Dictionary with:
                - text: Extracted text content
                - paragraph_count: Number of paragraphs
                - metadata: Document metadata
        """
        try:
            from docx import Document

            # Open DOCX from bytes
            doc = Document(BytesIO(docx_bytes))

            # Extract text from paragraphs
            paragraphs = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs.append(text)

            # Extract text from tables
            table_texts = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        table_texts.append(row_text)

            full_text = "\n\n".join(paragraphs)
            if table_texts:
                full_text += "\n\n[Tables]\n" + "\n".join(table_texts)

            metadata = {
                'paragraph_count': len(paragraphs),
                'table_count': len(doc.tables),
                'author': doc.core_properties.author or '',
                'title': doc.core_properties.title or '',
                'subject': doc.core_properties.subject or ''
            }

            logger.info(f"Extracted {len(full_text)} characters from DOCX with {metadata['paragraph_count']} paragraphs")

            return {
                'text': full_text,
                'paragraph_count': metadata['paragraph_count'],
                'metadata': metadata
            }

        except Exception as e:
            logger.error(f"Error extracting DOCX: {str(e)}", exc_info=True)
            return {
                'text': '',
                'paragraph_count': 0,
                'metadata': {},
                'error': str(e)
            }


class XLSXExtractor:
    """Extract text from Excel files using openpyxl."""

    def extract(self, xlsx_bytes: bytes) -> Dict[str, Any]:
        """
        Extract text from XLSX file.

        Args:
            xlsx_bytes: XLSX file content as bytes

        Returns:
            Dictionary with:
                - text: Extracted text content
                - sheet_data: Information about sheets
                - metadata: Workbook metadata
        """
        try:
            from openpyxl import load_workbook

            # Open XLSX from bytes
            wb = load_workbook(BytesIO(xlsx_bytes), read_only=True, data_only=True)

            sheet_texts = []
            sheet_data = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]

                # Extract cell values
                rows_text = []
                row_count = 0
                for row in sheet.iter_rows(values_only=True):
                    row_count += 1
                    # Filter out None values and convert to string
                    row_values = [str(cell) for cell in row if cell is not None and str(cell).strip()]
                    if row_values:
                        rows_text.append(' | '.join(row_values))

                sheet_text = f"[Sheet: {sheet_name}]\n" + "\n".join(rows_text)
                sheet_texts.append(sheet_text)

                sheet_data.append({
                    'name': sheet_name,
                    'row_count': row_count
                })

            full_text = "\n\n".join(sheet_texts)

            metadata = {
                'sheet_count': len(wb.sheetnames),
                'sheet_names': wb.sheetnames
            }

            wb.close()

            logger.info(f"Extracted {len(full_text)} characters from XLSX with {metadata['sheet_count']} sheets")

            return {
                'text': full_text,
                'sheet_data': sheet_data,
                'metadata': metadata
            }

        except Exception as e:
            logger.error(f"Error extracting XLSX: {str(e)}", exc_info=True)
            return {
                'text': '',
                'sheet_data': [],
                'metadata': {},
                'error': str(e)
            }


class PPTXExtractor:
    """Extract text from PowerPoint files using python-pptx."""

    def extract(self, pptx_bytes: bytes) -> Dict[str, Any]:
        """
        Extract text from PPTX file.

        Args:
            pptx_bytes: PPTX file content as bytes

        Returns:
            Dictionary with:
                - text: Extracted text content
                - slide_count: Number of slides
                - metadata: Presentation metadata
        """
        try:
            from pptx import Presentation

            # Open PPTX from bytes
            prs = Presentation(BytesIO(pptx_bytes))

            slide_texts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_parts = [f"[Slide {slide_num}]"]

                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_parts.append(shape.text.strip())

                slide_text = "\n".join(slide_parts)
                slide_texts.append(slide_text)

            full_text = "\n\n".join(slide_texts)

            metadata = {
                'slide_count': len(prs.slides),
                'author': prs.core_properties.author or '',
                'title': prs.core_properties.title or '',
                'subject': prs.core_properties.subject or ''
            }

            logger.info(f"Extracted {len(full_text)} characters from PPTX with {metadata['slide_count']} slides")

            return {
                'text': full_text,
                'slide_count': metadata['slide_count'],
                'metadata': metadata
            }

        except Exception as e:
            logger.error(f"Error extracting PPTX: {str(e)}", exc_info=True)
            return {
                'text': '',
                'slide_count': 0,
                'metadata': {},
                'error': str(e)
            }


class MSGExtractor:
    """Extract text from Outlook .msg files using extract-msg, or from embedded message text."""

    def extract(self, msg_bytes: bytes) -> Dict[str, Any]:
        """
        Extract text from .msg file or embedded message text.

        Args:
            msg_bytes: .msg file content as bytes, or formatted email text

        Returns:
            Dictionary with:
                - text: Extracted text content
                - metadata: Email metadata (subject, sender, recipients, etc.)
        """
        try:
            # First, try to decode as UTF-8 text (for embedded messages from Graph API)
            try:
                text_content = msg_bytes.decode('utf-8')
                # If it's already formatted text (from Graph API itemAttachment), use it directly
                if text_content.startswith('Subject:') or 'From:' in text_content:
                    logger.info("Processing embedded message text (not .msg file)")
                    # Parse the formatted text
                    lines = text_content.split('\n')
                    subject = ''
                    sender = ''
                    date = ''
                    body_start = 0
                    
                    for i, line in enumerate(lines):
                        if line.startswith('Subject:'):
                            subject = line.replace('Subject:', '').strip()
                        elif line.startswith('From:'):
                            sender = line.replace('From:', '').strip()
                        elif line.startswith('Date:'):
                            date = line.replace('Date:', '').strip()
                        elif line.strip() == '' and i > 0:
                            body_start = i + 1
                            break
                    
                    body = '\n'.join(lines[body_start:]) if body_start < len(lines) else text_content
                    
                    metadata = {
                        'subject': subject,
                        'sender': sender,
                        'date': date,
                        'is_embedded_message': True
                    }
                    
                    return {
                        'text': text_content,
                        'metadata': metadata
                    }
            except UnicodeDecodeError:
                # Not UTF-8 text, try as .msg file
                pass
            
            # Try to process as .msg file using extract-msg
            import extract_msg
            from io import BytesIO
            import tempfile
            import os

            # extract-msg requires a file path, so we need to write to temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.msg') as tmp_file:
                tmp_file.write(msg_bytes)
                tmp_path = tmp_file.name

            try:
                # Open the .msg file
                msg = extract_msg.Message(tmp_path)

                # Extract email content
                text_parts = []

                # Email metadata
                subject = msg.subject or ''
                sender = msg.sender or ''
                recipients = msg.to or ''
                cc = msg.cc or ''
                date = str(msg.date) if hasattr(msg, 'date') and msg.date else ''
                body = msg.body or ''
                html_body = msg.htmlBody or ''

                # Build structured text
                if subject:
                    text_parts.append(f"Subject: {subject}")
                if sender:
                    text_parts.append(f"From: {sender}")
                if recipients:
                    text_parts.append(f"To: {recipients}")
                if cc:
                    text_parts.append(f"CC: {cc}")
                if date:
                    text_parts.append(f"Date: {date}")
                
                text_parts.append("")  # Separator
                
                # Use HTML body if available (usually more complete), otherwise plain text
                if html_body:
                    # Strip HTML tags for cleaner text
                    try:
                        from bs4 import BeautifulSoup
                        soup = BeautifulSoup(html_body, 'html.parser')
                        body_text = soup.get_text(separator='\n', strip=True)
                        text_parts.append(body_text)
                    except Exception:
                        # Fallback to raw HTML if BeautifulSoup fails
                        text_parts.append(html_body)
                elif body:
                    text_parts.append(body)

                full_text = "\n".join(text_parts)

                # Get attachments info if any
                attachments_info = []
                if hasattr(msg, 'attachments') and msg.attachments:
                    for att in msg.attachments:
                        attachments_info.append({
                            'name': att.longFilename or att.shortFilename or 'unknown',
                            'size': getattr(att, 'dataLength', 0)
                        })

                metadata = {
                    'subject': subject,
                    'sender': sender,
                    'recipients': recipients,
                    'cc': cc,
                    'date': date,
                    'has_html': bool(html_body),
                    'attachment_count': len(attachments_info),
                    'attachments': attachments_info
                }

                msg.close()

                logger.info(f"Extracted {len(full_text)} characters from .msg file (Subject: {subject[:50]})")

                return {
                    'text': full_text,
                    'metadata': metadata
                }

            finally:
                # Clean up temp file
                try:
                    os.unlink(tmp_path)
                except Exception:
                    pass

        except Exception as e:
            logger.error(f"Error extracting .msg file: {str(e)}", exc_info=True)
            return {
                'text': '',
                'metadata': {},
                'error': str(e)
            }


class TextExtractor:
    """Extract text from plain text files."""

    def extract(self, text_bytes: bytes, encoding: str = 'utf-8') -> Dict[str, Any]:
        """
        Extract text from plain text file.

        Args:
            text_bytes: Text file content as bytes
            encoding: Text encoding (default: utf-8)

        Returns:
            Dictionary with:
                - text: Extracted text content
                - line_count: Number of lines
                - metadata: Text file metadata
        """
        try:
            # Try specified encoding first
            try:
                text = text_bytes.decode(encoding)
            except UnicodeDecodeError:
                # Fallback to chardet for encoding detection
                try:
                    import chardet
                    detected = chardet.detect(text_bytes)
                    encoding = detected['encoding'] or 'utf-8'
                    text = text_bytes.decode(encoding, errors='replace')
                    logger.info(f"Detected encoding: {encoding}")
                except Exception:
                    # Last resort: decode with errors='replace'
                    text = text_bytes.decode('utf-8', errors='replace')
                    logger.warning("Using UTF-8 with error replacement")

            lines = text.split('\n')
            line_count = len(lines)

            metadata = {
                'line_count': line_count,
                'encoding': encoding,
                'size_bytes': len(text_bytes)
            }

            logger.info(f"Extracted {len(text)} characters from text file with {line_count} lines")

            return {
                'text': text,
                'line_count': line_count,
                'metadata': metadata
            }

        except Exception as e:
            logger.error(f"Error extracting text: {str(e)}", exc_info=True)
            return {
                'text': '',
                'line_count': 0,
                'metadata': {},
                'error': str(e)
            }


class DocumentExtractorFactory:
    """Factory to select appropriate extractor based on MIME type."""

    MIME_TO_EXTRACTOR = {
        'application/pdf': PDFExtractor,
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': DOCXExtractor,
        'application/msword': DOCXExtractor,  # Legacy .doc (will attempt with python-docx)
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': XLSXExtractor,
        'application/vnd.ms-excel': XLSXExtractor,  # Legacy .xls
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': PPTXExtractor,
        'application/vnd.ms-powerpoint': PPTXExtractor,  # Legacy .ppt
        'application/vnd.ms-outlook': MSGExtractor,  # Outlook .msg files
        'application/x-msg': MSGExtractor,  # Alternative MIME type for .msg
        'message/rfc822': MSGExtractor,  # RFC 822 message format
        'text/plain': TextExtractor,
        'text/csv': TextExtractor,
        'text/html': TextExtractor,
        'text/xml': TextExtractor,
        'application/json': TextExtractor
    }

    @staticmethod
    def get_extractor(mime_type: str, filename: str = '') -> Optional[Any]:
        """
        Get appropriate extractor for given MIME type or filename.

        Args:
            mime_type: MIME type of the file
            filename: Optional filename to detect type by extension

        Returns:
            Extractor class instance or None if not supported
        """
        extractor_class = DocumentExtractorFactory.MIME_TO_EXTRACTOR.get(mime_type)
        
        # If no MIME type match, try to detect by file extension
        if not extractor_class and filename:
            filename_lower = filename.lower()
            if filename_lower.endswith('.msg'):
                extractor_class = MSGExtractor
            elif filename_lower.endswith('.pdf'):
                extractor_class = PDFExtractor
            elif filename_lower.endswith(('.docx', '.doc')):
                extractor_class = DOCXExtractor
            elif filename_lower.endswith(('.xlsx', '.xls')):
                extractor_class = XLSXExtractor
            elif filename_lower.endswith(('.pptx', '.ppt')):
                extractor_class = PPTXExtractor
            elif filename_lower.endswith(('.txt', '.csv', '.html', '.xml', '.json')):
                extractor_class = TextExtractor
        
        if extractor_class:
            return extractor_class()
        else:
            logger.warning(f"No extractor found for MIME type: {mime_type}, filename: {filename}")
            return None

    @staticmethod
    def is_supported(mime_type: str, filename: str = '') -> bool:
        """
        Check if MIME type or filename extension is supported.

        Args:
            mime_type: MIME type to check
            filename: Optional filename to check extension

        Returns:
            bool: True if supported
        """
        if mime_type in DocumentExtractorFactory.MIME_TO_EXTRACTOR:
            return True
        
        # Check by file extension if MIME type not found
        if filename:
            filename_lower = filename.lower()
            if filename_lower.endswith(('.msg', '.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.txt', '.csv', '.html', '.xml', '.json')):
                return True
        
        return False

    @staticmethod
    def get_supported_types() -> list:
        """
        Get list of supported MIME types.

        Returns:
            List of supported MIME types
        """
        return list(DocumentExtractorFactory.MIME_TO_EXTRACTOR.keys())
